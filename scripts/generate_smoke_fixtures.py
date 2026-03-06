from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from markitdown import MarkItDown
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1] / "tests" / "fixtures" / "smoke"


def build_docx(path: Path) -> None:
    document = Document()
    document.add_heading("Sample DOCX", level=1)
    document.add_paragraph("Hello from docx fixture.")
    document.save(path)


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    presentation.save(path)


def build_xlsx(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet["A1"] = "Name"
    worksheet["B1"] = "Score"
    worksheet.append(["Alice", 42])
    workbook.save(path)


def build_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path), pagesize=letter)
    pdf.drawString(72, 720, "Sample PDF")
    pdf.drawString(72, 700, "Hello from pdf fixture.")
    pdf.save()


def main() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)

    build_docx(ROOT / "sample.docx")
    build_pptx(ROOT / "sample.pptx")
    build_xlsx(ROOT / "sample.xlsx")
    build_pdf(ROOT / "sample.pdf")
    (ROOT / "broken.docx").write_bytes(bytes(range(256)))

    converter = MarkItDown(enable_plugins=False)
    expected_markdown = {}
    for filename in ("sample.docx", "sample.pptx", "sample.xlsx", "sample.pdf"):
        result = converter.convert(str(ROOT / filename))
        expected_markdown[filename] = result.text_content.replace("\r\n", "\n").strip()

    (ROOT / "expected-markdown.json").write_text(
        json.dumps(expected_markdown, indent=2) + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
