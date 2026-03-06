from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def build_sample_fixtures(root: Path) -> None:
    root.mkdir(parents=True, exist_ok=True)

    docx_path = root / "sample.docx"
    document = Document()
    document.add_heading("Sample DOCX", level=1)
    document.add_paragraph("Hello from docx fixture.")
    document.save(docx_path)

    pptx_path = root / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    presentation.save(pptx_path)

    xlsx_path = root / "sample.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet["A1"] = "Name"
    worksheet["B1"] = "Score"
    worksheet.append(["Alice", 42])
    workbook.save(xlsx_path)

    pdf_path = root / "sample.pdf"
    pdf = canvas.Canvas(str(pdf_path), pagesize=letter)
    pdf.drawString(72, 720, "Sample PDF")
    pdf.drawString(72, 700, "Hello from pdf fixture.")
    pdf.save()

    (root / "broken.docx").write_bytes(bytes(range(256)))
